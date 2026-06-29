import os
import json
import urllib.parse
from http.server import BaseHTTPRequestHandler, HTTPServer
import pandas as pd
import pmo_weekly_report
import pdfplumber
import re
import math
import threading
from supabase import create_client

SUPABASE_URL = "https://xczrfhxlzfrytbtpidhn.supabase.co"
SUPABASE_KEY = "sb_publishable_wngiEluCpegeAyfehAkQ0Q_0P6Fmklw"
supabase = create_client(SUPABASE_URL, SUPABASE_KEY)

EXCEL_LOCK = threading.Lock()

def sanitize_nan(val):
    if isinstance(val, dict):
        return {k: sanitize_nan(v) for k, v in val.items()}
    elif isinstance(val, list):
        return [sanitize_nan(v) for v in val]
    try:
        # Check if it is a pandas NA or null value
        if pd.isna(val):
            return None
    except:
        pass
    if isinstance(val, float) and (math.isnan(val) or math.isinf(val)):
        return None
    return val


PORT = int(os.environ.get("PORT", 8000))
EXCEL_PATH = "PMO_Weekly_Report_Data.xlsx"

# ---------------------------------------------------------------------------
# PDF GANTT SCHEDULE PARSER
# ---------------------------------------------------------------------------
def clean_cell(text):
    if not text:
        return ""
    return re.sub(r'\s+', ' ', text.strip())

def parse_pdf_plan_with_hierarchy(path):
    with pdfplumber.open(path) as pdf:
        tasks = []
        left_pages = [0, 1, 2, 3]
        right_pages = [4, 5, 6, 7]
        
        for idx in range(4):
            if idx >= len(pdf.pages):
                break
            lp = pdf.pages[left_pages[idx]]
            
            # Use right page mapping if available, otherwise fallback
            rp_idx = right_pages[idx]
            if rp_idx < len(pdf.pages):
                rp = pdf.pages[rp_idx]
                r_tables = rp.extract_tables()
            else:
                r_tables = []
                
            l_tables = lp.extract_tables()
            
            if not l_tables:
                continue
                
            l_rows = l_tables[0]
            r_rows = r_tables[0] if r_tables else []
            
            words = lp.extract_words()
            
            for r_idx in range(1, len(l_rows)):
                l_row = l_rows[r_idx]
                r_row = r_rows[r_idx] if r_idx < len(r_rows) else [""] * 8
                
                tid = clean_cell(l_row[0])
                if not tid:
                    continue
                    
                name = clean_cell(l_row[3])
                
                # Check coordinates of first word
                first_word = re.findall(r'[A-Za-z0-9\-]+', name)
                x0_val = 161.7  # fallback to level 3
                
                if first_word:
                    target_word = first_word[0].lower()
                    matching_words = [w for w in words if w["text"].lower().startswith(target_word) and 100 < w["x0"] < 250]
                    if matching_words:
                        estimated_top = 50 + r_idx * 15.6
                        matching_words.sort(key=lambda w: abs(w["top"] - estimated_top))
                        x0_val = matching_words[0]["x0"]
                
                # Deduce hierarchy level using the baseline-offset math formula
                level = int(round((x0_val - 139.1) / 11.25)) + 1
                if level < 1:
                    level = 1
                    
                task = {
                    "id": int(tid),
                    "name": name,
                    "duration": clean_cell(l_row[4]),
                    "baselineStart": clean_cell(l_row[5]),
                    "baselineFinish": clean_cell(l_row[6]),
                    "start": clean_cell(l_row[7]),
                    "finish": clean_cell(r_row[0]) if len(r_row) > 0 else "",
                    "percentComplete": clean_cell(r_row[1]) if len(r_row) > 1 else "",
                    "progressPercent": clean_cell(r_row[2]) if len(r_row) > 2 else "",
                    "resources": clean_cell(r_row[3]) if len(r_row) > 3 else "",
                    "predecessors": clean_cell(r_row[4]) if len(r_row) > 4 else "",
                    "comments": clean_cell(r_row[5]) if len(r_row) > 5 else "",
                    "deliverable": clean_cell(r_row[6]) if len(r_row) > 6 else "",
                    "level": level
                }
                tasks.append(task)
                
        tasks.sort(key=lambda x: x["id"])
        
        project_name = "Freedom Telecom"
        if tasks:
            project_name = tasks[0]["name"]
            
        return {
            "projectName": project_name,
            "tasks": tasks
        }


def read_document_text(file_path, filename):
    text = ""
    ext = os.path.splitext(filename)[1].lower()
    
    if ext == ".pdf":
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        text += t + "\n"
        except Exception as e:
            print(f"Error extracting text from PDF: {e}")
    elif ext == ".docx":
        try:
            import docx
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        except Exception as e:
            print(f"Error extracting text from DOCX: {e}")
    else: # Fallback to txt
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        except Exception as e:
            print(f"Error reading TXT file: {e}")
    return text


def extract_scope_and_components(file_path, filename):
    text = read_document_text(file_path, filename)

    # Heuristics for components list
    common_components = [
        "Billing", "Provisioning", "CRM", "SMS", "USSD", "Gateway", "Portal", 
        "Database", "Reporting", "Analytics", "Dashboard", "Charging", "BSS", 
        "OSS", "Integration", "API", "Security", "HLR", "HSS", "PCRF", "VMS", 
        "IVR", "MMSC", "Notification", "E-Commerce", "Mobile App"
    ]
    
    found_components = []
    
    # Check for direct keyword matches in text
    for comp in common_components:
        pattern = r'\b' + re.escape(comp) + r'\b'
        if re.search(pattern, text, re.IGNORECASE):
            found_components.append(comp + " Module")

    lines = text.split("\n")
    collect_bullets = False
    bullet_patterns = [r'^\s*-\s*(.+)', r'^\s*\*\s*(.+)', r'^\s*•\s*(.+)', r'^\s*\d+\.\s*(.+)']
    
    for line in lines:
        cleaned = line.strip()
        if not cleaned:
            continue
        
        lower_line = cleaned.lower()
        if any(h in lower_line for h in ["components list", "modules included", "scope items", "systems list", "key modules"]):
            collect_bullets = True
            continue
        elif collect_bullets and cleaned.endswith(":") and len(cleaned) < 30:
            collect_bullets = False
            
        if collect_bullets:
            for pat in bullet_patterns:
                m = re.match(pat, cleaned)
                if m:
                    bullet_text = m.group(1).strip()
                    if 3 < len(bullet_text) < 50:
                        if bullet_text not in found_components:
                            found_components.append(bullet_text)
                    break
                    
    # Clean up and ensure unique components list
    if not found_components:
        found_components = ["Core Engine", "Database Integration", "User Authentication", "Web Dashboard Portal"]
    else:
        unique_list = []
        for x in found_components:
            if x not in unique_list:
                unique_list.append(x)
        found_components = unique_list[:12] # Limit to top 12

    # Heuristic for Project Manager
    pm_name = "Madhava Reddy Gorreegala"
    pm_match = re.search(r'(?:Project|Program)\s+Manager\s*[:\-]?\s*([^\n\r]+)', text, re.IGNORECASE)
    if pm_match:
        pm_name = pm_match.group(1).strip()
    else:
        for line in lines:
            if "manager" in line.lower() and len(line.strip()) < 50:
                parts = re.split(r'[:\-]', line)
                if len(parts) > 1:
                    pm_name = parts[-1].strip()
                    break

    # Heuristic for Latest Updates
    updates = []
    collect_updates = False
    for line in lines:
        cleaned = line.strip()
        if not cleaned:
            continue
        if any(h in cleaned.lower() for h in ["latest updates", "progress update", "weekly updates", "recent updates", "accomplishments"]):
            collect_updates = True
            continue
        elif collect_updates and (cleaned.endswith(":") or any(h in cleaned.lower() for h in ["risk", "milestone", "support", "planned activities"])):
            collect_updates = False
            
        if collect_updates:
            for pat in bullet_patterns:
                m = re.match(pat, cleaned)
                if m:
                    updates.append(m.group(1).strip())
                    break
            else:
                if len(cleaned) > 15:
                    updates.append(cleaned)
    if not updates:
        updates = [
            "Customer onboarding via the Mobile App has been successfully completed",
            "Airalo (eSIM) integration and testing have been completed",
            "Updated APIs shared with the client for integration",
            "All core BSS services are working fine"
        ]

    # Heuristic for Risks & Dependencies
    risks = []
    collect_risks = False
    for line in lines:
        cleaned = line.strip()
        if not cleaned:
            continue
        if any(h in cleaned.lower() for h in ["risk and dependencies", "risks and dependencies", "key risks", "risk register", "risk & dependencies"]):
            collect_risks = True
            continue
        elif collect_risks and (cleaned.endswith(":") or any(h in cleaned.lower() for h in ["latest updates", "milestone", "support", "planned"])):
            collect_risks = False
            
        if collect_risks:
            for pat in bullet_patterns:
                m = re.match(pat, cleaned)
                if m:
                    risks.append(m.group(1).strip())
                    break
            else:
                if len(cleaned) > 15:
                    risks.append(cleaned)
    if not risks:
        risks = ["Changes/additions in the requirements are delaying the production integrations"]

    # Heuristic for Support Required
    support = []
    collect_support = False
    for line in lines:
        cleaned = line.strip()
        if not cleaned:
            continue
        if any(h in cleaned.lower() for h in ["support required", "help needed", "escalations", "support request"]):
            collect_support = True
            continue
        elif collect_support and (cleaned.endswith(":") or any(h in cleaned.lower() for h in ["latest updates", "milestone", "risk", "planned"])):
            collect_support = False
            
        if collect_support:
            for pat in bullet_patterns:
                m = re.match(pat, cleaned)
                if m:
                    support.append(m.group(1).strip())
                    break
            else:
                if len(cleaned) > 15:
                    support.append(cleaned)
    if not support:
        support = [
            "Dealer Mobile App testing completion",
            "MNP process to be deployed in production environment",
            "AWS Infra security configuration and port validation"
        ]

    milestones = [
        {"Activity": "Kick-off meeting", "ETA": "Mar-24", "Rev": "", "UpdatedPlan": "", "Status": "Completed", "Comments": ""},
        {"Activity": "BSS Demo workshop", "ETA": "Mar-24", "Rev": "", "UpdatedPlan": "", "Status": "Completed", "Comments": ""},
        {"Activity": "Input data from Liv.ing", "ETA": "May-24", "Rev": "", "UpdatedPlan": "", "Status": "Completed", "Comments": ""},
        {"Activity": "SOW preparation & Sign off", "ETA": "Jul-24", "Rev": "", "UpdatedPlan": "", "Status": "Completed", "Comments": ""},
        {"Activity": "AWS Infra Deployment (Testing)", "ETA": "Jul-24", "Rev": "", "UpdatedPlan": "", "Status": "Completed", "Comments": ""},
        {"Activity": "AWS Production re-deployment", "ETA": "Oct-24", "Rev": "Dec-25", "UpdatedPlan": "May-25", "Status": "On Track", "Comments": ""},
        {"Activity": "Core integration", "ETA": "Oct-24", "Rev": "Jan-25", "UpdatedPlan": "Oct-25", "Status": "On Track", "Comments": ""},
        {"Activity": "E2E API - Integrations", "ETA": "Oct-24", "Rev": "Jan-25", "UpdatedPlan": "Oct-25", "Status": "On Track", "Comments": ""},
        {"Activity": "Integration and Testing", "ETA": "Oct-24", "Rev": "Jan-25", "UpdatedPlan": "Dec-25", "Status": "At Risk", "Comments": ""},
        {"Activity": "UAT", "ETA": "Oct-24", "Rev": "Jan-25", "UpdatedPlan": "Dec-25", "Status": "At Risk", "Comments": ""},
        {"Activity": "Go-live", "ETA": "Oct-24", "Rev": "Jan-25", "UpdatedPlan": "Dec-25", "Status": "At Risk", "Comments": ""},
        {"Activity": "Training and Handover", "ETA": "Oct-24", "Rev": "Jan-25", "UpdatedPlan": "Dec-25", "Status": "At Risk", "Comments": ""},
    ]

    for row in milestones:
        act_escaped = re.escape(row["Activity"])
        match = re.search(act_escaped + r'\s*\|\s*([A-Za-z0-9\-]+)', text, re.IGNORECASE)
        if match:
            row["ETA"] = match.group(1).strip()

    # Formulate parsed scope fields
    scope_fields = {
        "Pre-sales Document": filename,
        "Scope Prepare": "Extracted from " + filename,
        "Project Time Plan": "Drafted based on " + filename,
        "Project Status": f"Pre-sales phase initiated. Detected {len(found_components)} project modules.",
        "Project Manager": pm_name,
        "Latest Updates": "\n".join(updates),
        "Risks Dependencies": "\n".join(risks),
        "Support Required": "\n".join(support),
        "Milestones": json.dumps(milestones)
    }

    return scope_fields, found_components


def extract_agreed_document_details(file_path, filename):
    text = read_document_text(file_path, filename)
    
    project_plan = ""
    scope_in = ""
    scope_out = ""
    
    lines = text.split("\n")
    current_section = None
    
    for line in lines:
        cleaned = line.strip()
        if not cleaned:
            continue
        
        lower_line = cleaned.lower()
        # Section headers detection
        if any(h in lower_line for h in ["project plan", "project timeline", "milestones timeline", "schedule"]):
            current_section = "plan"
            continue
        elif any(h in lower_line for h in ["scope in", "in scope", "deliverables included", "inclusions"]):
            current_section = "in"
            continue
        elif any(h in lower_line for h in ["scope out", "out of scope", "exclusions", "not included"]):
            current_section = "out"
            continue
        elif cleaned.endswith(":") and len(cleaned) < 25:
            current_section = None
            continue
            
        if current_section == "plan":
            project_plan += cleaned + "\n"
        elif current_section == "in":
            scope_in += cleaned + "\n"
        elif current_section == "out":
            scope_out += cleaned + "\n"
            
    # Fallbacks if sections not explicitly found
    if not project_plan.strip():
        time_lines = [l.strip() for l in lines if any(k in l.lower() for k in ["week", "month", "milestone", "date", "phase"])]
        project_plan = "\n".join(time_lines[:5]) if time_lines else "Project baseline plan drafted based on " + filename
    if not scope_in.strip():
        in_lines = [l.strip() for l in lines if any(k in l.lower() for k in ["deliver", "require", "crm", "billing", "system", "portal"])]
        scope_in = "\n".join(in_lines[:5]) if in_lines else "Standard project delivery scope."
    if not scope_out.strip():
        out_lines = [l.strip() for l in lines if any(k in l.lower() for k in ["exclude", "third-party", "customization", "not responsible"])]
        scope_out = "\n".join(out_lines[:5]) if out_lines else "Any item not explicitly stated in the In-Scope section is considered out of scope."
        
    return {
        "Agreed Document": filename,
        "Project Plan": project_plan.strip()[:1000],
        "Scope In": scope_in.strip()[:1000],
        "Scope Out": scope_out.strip()[:1000]
    }


class PMORequestHandler(BaseHTTPRequestHandler):
    def log_message(self, format, *args):
        # Silence default logging to keep terminal output clean
        pass

    def check_auth(self):
        auth_header = self.headers.get("Authorization", "")
        if not auth_header.startswith("Bearer "):
            self.send_error_response(401, "Missing or invalid authorization token.")
            return False
        token = auth_header.split("Bearer ")[1].strip()
        try:
            res = supabase.auth.get_user(token)
            if res and res.user:
                return True
        except Exception as e:
            print(f"Token validation failed: {e}")
        self.send_error_response(401, "Session expired or invalid token.")
        return False

    def do_GET(self):
        parsed_path = urllib.parse.urlparse(self.path)
        path = parsed_path.path

        # Route: API Data Retrieval
        if path == "/api/data":
            if not self.check_auth():
                return

            try:
                # Query projects
                proj_res = supabase.table("projects").select("*").execute()
                projects = []
                for p in (proj_res.data or []):
                    projects.append({
                        "Project": p.get("project"),
                        "Manager": p.get("manager"),
                        "Planned %": p.get("planned_pct"),
                        "Actual %": p.get("actual_pct"),
                        "Budget": p.get("budget"),
                        "Actual Cost": p.get("actual_cost"),
                        "Delay Days": p.get("delay_days"),
                        "Risk Score": p.get("risk_score"),
                        "Approved By": p.get("approved_by")
                    })

                # Query achievements
                ach_res = supabase.table("achievements").select("*").execute()
                achievements = []
                for a in (ach_res.data or []):
                    achievements.append({
                        "Project": a.get("project"),
                        "Achievement": a.get("achievement")
                    })

                # Query risks
                risk_res = supabase.table("risks").select("*").execute()
                risks = []
                for r in (risk_res.data or []):
                    risks.append({
                        "Project": r.get("project"),
                        "Risk Description": r.get("risk_description"),
                        "Probability": r.get("probability"),
                        "Impact": r.get("impact"),
                        "Mitigation": r.get("mitigation")
                    })

                # Query nextsteps
                ns_res = supabase.table("next_steps").select("*").execute()
                nextsteps = []
                for n in (ns_res.data or []):
                    nextsteps.append({
                        "Project": n.get("project"),
                        "Task": n.get("task"),
                        "Owner": n.get("owner"),
                        "Deadline": n.get("deadline")
                    })

                # Query decisions
                dec_res = supabase.table("decisions").select("*").execute()
                decisions = []
                for d in (dec_res.data or []):
                    decisions.append({
                        "Project": d.get("project"),
                        "Decision Required": d.get("decision_required"),
                        "Context": d.get("context"),
                        "Options": d.get("options"),
                        "Recommendation": d.get("recommendation")
                    })

                # Query scopes
                scope_res = supabase.table("project_scopes").select("*").execute()
                scopes = []
                for s in (scope_res.data or []):
                    scopes.append({
                        "Project": s.get("project"),
                        "Pre-sales Document": s.get("presales_document"),
                        "Components": s.get("components"),
                        "Scope Prepare": s.get("scope_prepare"),
                        "Project Time Plan": s.get("project_time_plan"),
                        "Agreed Document": s.get("agreed_document"),
                        "Project Plan": s.get("project_plan"),
                        "Scope In": s.get("scope_in"),
                        "Scope Out": s.get("scope_out"),
                        "Customer Presentation Type": s.get("customer_presentation_type"),
                        "Kick Off": s.get("kick_off"),
                        "Weekly": s.get("weekly"),
                        "Monthly": s.get("monthly"),
                        "Milestone Progress": s.get("milestone_progress"),
                        "Project Status": s.get("project_status"),
                        "Monthly Status": s.get("monthly_status"),
                        "Executive Presentation": s.get("executive_presentation"),
                        "Risk Register Party Expected": s.get("risk_register_party_expected")
                    })

                # Query latestApproval
                app_res = supabase.table("approvals").select("*").order("timestamp", desc=True).limit(1).execute()
                latest_approval = None
                if app_res.data and len(app_res.data) > 0:
                    app = app_res.data[0]
                    latest_approval = {
                        "Timestamp": app.get("timestamp"),
                        "Approved By": app.get("approved_by"),
                        "Approved Projects": app.get("approved_projects")
                    }

                data = {
                    "projects": projects,
                    "achievements": achievements,
                    "risks": risks,
                    "nextsteps": nextsteps,
                    "decisions": decisions,
                    "scopes": scopes,
                    "latestApproval": latest_approval
                }

                # Format dates in NextSteps to strings for JSON safety
                for item in data["nextsteps"]:
                    if "Deadline" in item and item["Deadline"]:
                        item["Deadline"] = str(item["Deadline"])

                self.send_json_response(200, data)
            except Exception as e:
                import traceback
                traceback.print_exc()
                self.send_error_response(500, f"Error querying Supabase database: {str(e)}")
            return

        elif path == "/api/plan/data":
            if not self.check_auth():
                return
            plan_path = "project_plan.json"
            if not os.path.exists(plan_path):
                self.send_json_response(200, {"projectName": "", "tasks": []})
                return
            try:
                with open(plan_path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                self.send_json_response(200, data)
            except Exception as e:
                self.send_error_response(500, f"Error reading plan: {str(e)}")
            return

        # Route: Static File Serving
        # Map URL path to local file path
        if path == "/" or path == "/index.html":
            local_file = "index.html"
            content_type = "text/html"
            is_binary = False
        elif path == "/styles.css":
            local_file = "styles.css"
            content_type = "text/css"
            is_binary = False
        elif path == "/app.js":
            local_file = "app.js"
            content_type = "application/javascript"
            is_binary = False
        elif path == "/logo.png":
            local_file = "logo.png"
            content_type = "image/png"
            is_binary = True
        elif path == "/api/download/docx":
            local_file = "PMO_Executive_Report.docx"
            content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            is_binary = True
        elif path == "/api/download/excel":
            local_file = EXCEL_PATH
            content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            is_binary = True
        else:
            self.send_error_response(404, "File not found.")
            return

        # Check if static file exists and serve it
        if os.path.exists(local_file):
            try:
                if is_binary:
                    with open(local_file, "rb") as f:
                        content = f.read()
                else:
                    with open(local_file, "r", encoding="utf-8") as f:
                        content = f.read().encode('utf-8')
                        
                self.send_response(200)
                self.send_header("Content-Type", content_type)
                self.send_header("Content-Length", str(len(content)))
                # Disable caching for instant updates
                self.send_header("Cache-Control", "no-store, no-cache, must-revalidate, max-age=0")
                self.send_header("Pragma", "no-cache")
                self.send_header("Expires", "0")
                # Add download header for API download routes
                if path.startswith("/api/download/"):
                    filename = os.path.basename(local_file)
                    self.send_header("Content-Disposition", f'attachment; filename="{filename}"')
                self.end_headers()
                self.wfile.write(content)
            except Exception as e:
                self.send_error_response(500, f"Error serving file: {str(e)}")
        else:
            self.send_error_response(404, f"File {local_file} not found.")

    def do_POST(self):
        parsed_path = urllib.parse.urlparse(self.path)
        path = parsed_path.path

        # Check token authorization for all POST APIs
        if path.startswith("/api/"):
            if not self.check_auth():
                return

        # Route: Submit Data from Web UI
        if path == "/api/submit":
            try:
                content_length = int(self.headers['Content-Length'])
                post_data = self.rfile.read(content_length)
                payload = json.loads(post_data.decode('utf-8'))

                # Validate data structure
                required_keys = ["projects", "achievements", "risks", "nextsteps", "decisions"]
                if not all(k in payload for k in required_keys):
                    self.send_error_response(400, "Invalid payload. Missing sheet datasets.")
                    return

                # 1. Sync to Supabase
                # A. Upsert Projects
                for p in payload["projects"]:
                    supabase.table("projects").upsert({
                        "project": p.get("Project"),
                        "manager": p.get("Manager"),
                        "planned_pct": p.get("Planned %"),
                        "actual_pct": p.get("Actual %"),
                        "budget": p.get("Budget"),
                        "actual_cost": p.get("Actual Cost"),
                        "delay_days": p.get("Delay Days"),
                        "risk_score": p.get("Risk Score"),
                        "approved_by": p.get("Approved By")
                    }).execute()

                # B. Sync Deletions
                all_db_projects = supabase.table("projects").select("project").execute()
                db_proj_names = [p["project"] for p in (all_db_projects.data or [])]
                payload_proj_names = [p.get("Project") for p in payload["projects"]]
                deleted_proj_names = [name for name in db_proj_names if name not in payload_proj_names]
                for name in deleted_proj_names:
                    supabase.table("projects").delete().eq("project", name).execute()

                # C. Achievements
                if payload_proj_names:
                    supabase.table("achievements").delete().in_("project", payload_proj_names).execute()
                if "achievements" in payload and payload["achievements"]:
                    ach_rows = [{"project": a.get("Project"), "achievement": a.get("Achievement")} for a in payload["achievements"]]
                    supabase.table("achievements").insert(ach_rows).execute()

                # D. Risks
                if payload_proj_names:
                    supabase.table("risks").delete().in_("project", payload_proj_names).execute()
                if "risks" in payload and payload["risks"]:
                    risk_rows = [{
                        "project": r.get("Project"),
                        "risk_description": r.get("Risk Description"),
                        "probability": r.get("Probability"),
                        "impact": r.get("Impact"),
                        "mitigation": r.get("Mitigation")
                    } for r in payload["risks"]]
                    supabase.table("risks").insert(risk_rows).execute()

                # E. Next Steps
                if payload_proj_names:
                    supabase.table("next_steps").delete().in_("project", payload_proj_names).execute()
                if "nextsteps" in payload and payload["nextsteps"]:
                    ns_rows = [{
                        "project": n.get("Project"),
                        "task": n.get("Task"),
                        "owner": n.get("Owner"),
                        "deadline": n.get("Deadline") if n.get("Deadline") else None
                    } for n in payload["nextsteps"]]
                    supabase.table("next_steps").insert(ns_rows).execute()

                # F. Decisions
                if payload_proj_names:
                    supabase.table("decisions").delete().in_("project", payload_proj_names).execute()
                if "decisions" in payload and payload["decisions"]:
                    dec_rows = [{
                        "project": d.get("Project"),
                        "decision_required": d.get("Decision Required"),
                        "context": d.get("Context"),
                        "options": d.get("Options"),
                        "recommendation": d.get("Recommendation")
                    } for d in payload["decisions"]]
                    supabase.table("decisions").insert(dec_rows).execute()

                # G. Scopes
                if "scopes" in payload and payload["scopes"]:
                    for s in payload["scopes"]:
                        supabase.table("project_scopes").upsert({
                            "project": s.get("Project"),
                            "presales_document": s.get("Pre-sales Document"),
                            "components": s.get("Components"),
                            "scope_prepare": s.get("Scope Prepare"),
                            "project_time_plan": s.get("Project Time Plan"),
                            "agreed_document": s.get("Agreed Document"),
                            "project_plan": s.get("Project Plan"),
                            "scope_in": s.get("Scope In"),
                            "scope_out": s.get("Scope Out"),
                            "customer_presentation_type": s.get("Customer Presentation Type"),
                            "kick_off": s.get("Kick Off"),
                            "weekly": s.get("Weekly"),
                            "monthly": s.get("Monthly"),
                            "milestone_progress": s.get("Milestone Progress"),
                            "project_status": s.get("Project Status"),
                            "monthly_status": s.get("Monthly Status"),
                            "executive_presentation": s.get("Executive Presentation"),
                            "risk_register_party_expected": s.get("Risk Register Party Expected")
                        }).execute()

                # H. Approval Log history insertion
                approved_by = payload.get("approvedBy", "Project Lead")
                approved_projs_list = payload.get("approvedProjects", [])
                approved_projs_str = ", ".join(approved_projs_list) if approved_projs_list else ""
                if "approvedBy" in payload and payload["approvedBy"]:
                    supabase.table("approvals").insert({
                        "approved_by": approved_by,
                        "approved_projects": approved_projs_str
                    }).execute()

                # 2. Write to Excel compile cache
                df_projects = pd.DataFrame(payload["projects"])
                df_achievements = pd.DataFrame(payload["achievements"])
                df_risks = pd.DataFrame(payload["risks"])
                df_nextsteps = pd.DataFrame(payload["nextsteps"])
                df_decisions = pd.DataFrame(payload["decisions"])

                df_scopes = None
                if "scopes" in payload:
                    df_scopes = pd.DataFrame(payload["scopes"])
                elif os.path.exists(EXCEL_PATH):
                    try:
                        excel_file = pd.ExcelFile(EXCEL_PATH)
                        if "ProjectScope" in excel_file.sheet_names:
                            df_scopes = pd.read_excel(excel_file, sheet_name="ProjectScope")
                    except Exception as e:
                        print(f"Error reading existing ProjectScope: {e}")
                
                if df_scopes is None:
                    df_scopes = pd.DataFrame()

                # Align columns
                p_cols = ["Project", "Manager", "Planned %", "Actual %", "Budget", "Actual Cost", "Delay Days", "Risk Score", "Approved By"]
                for col in p_cols:
                    if col not in df_projects.columns:
                        df_projects[col] = None
                df_projects = df_projects[p_cols]

                a_cols = ["Project", "Achievement"]
                for col in a_cols:
                    if col not in df_achievements.columns:
                        df_achievements[col] = None
                df_achievements = df_achievements[a_cols]

                r_cols = ["Project", "Risk Description", "Probability", "Impact", "Mitigation"]
                for col in r_cols:
                    if col not in df_risks.columns:
                        df_risks[col] = None
                df_risks = df_risks[r_cols]

                n_cols = ["Project", "Task", "Owner", "Deadline"]
                for col in n_cols:
                    if col not in df_nextsteps.columns:
                        df_nextsteps[col] = None
                df_nextsteps = df_nextsteps[n_cols]

                d_cols = ["Project", "Decision Required", "Context", "Options", "Recommendation"]
                for col in d_cols:
                    if col not in df_decisions.columns:
                        df_decisions[col] = None
                df_decisions = df_decisions[d_cols]

                s_cols = [
                    "Project", "Pre-sales Document", "Components", "Scope Prepare", "Project Time Plan",
                    "Agreed Document", "Project Plan", "Scope In", "Scope Out", "Customer Presentation Type",
                    "Kick Off", "Weekly", "Monthly", "Milestone Progress", "Project Status", "Monthly Status",
                    "Executive Presentation", "Risk Register Party Expected"
                ]
                for col in s_cols:
                    if col not in df_scopes.columns:
                        df_scopes[col] = None
                df_scopes = df_scopes[s_cols]

                with EXCEL_LOCK:
                    df_approvals = pd.DataFrame(columns=["Timestamp", "Approved By", "Approved Projects"])
                    if os.path.exists(EXCEL_PATH):
                        try:
                            excel_file = pd.ExcelFile(EXCEL_PATH)
                            if "Approvals" in excel_file.sheet_names:
                                df_approvals = pd.read_excel(excel_file, sheet_name="Approvals")
                        except Exception as e:
                            print(f"Error reading approvals log: {e}")
                    
                    from datetime import datetime
                    new_row = {
                        "Timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                        "Approved By": approved_by,
                        "Approved Projects": approved_projs_str
                    }
                    df_approvals = pd.concat([df_approvals, pd.DataFrame([new_row])], ignore_index=True)

                    with pd.ExcelWriter(EXCEL_PATH, engine="openpyxl") as writer:
                        df_projects.to_excel(writer, sheet_name="Projects", index=False)
                        df_achievements.to_excel(writer, sheet_name="Achievements", index=False)
                        df_risks.to_excel(writer, sheet_name="Risks", index=False)
                        df_nextsteps.to_excel(writer, sheet_name="NextSteps", index=False)
                        df_decisions.to_excel(writer, sheet_name="Decisions", index=False)
                        df_scopes.to_excel(writer, sheet_name="ProjectScope", index=False)
                        df_approvals.to_excel(writer, sheet_name="Approvals", index=False)

                    print(f"Excel cache updated (Approved By: {approved_by}). Running report generation...")
                    pmo_weekly_report.run_report_generation(approved_by=approved_by)

                self.send_json_response(200, {
                    "status": "success",
                    "message": "Project data successfully synced with database and Word report generated."
                })

            except Exception as e:
                import traceback
                print(f"Error during submit: {e}")
                traceback.print_exc()
                self.send_error_response(500, f"Error processing submission: {str(e)}")
            return

        elif path == "/api/scope/upload":
            try:
                content_type = self.headers.get('Content-Type', '')
                if not content_type.startswith('multipart/form-data'):
                    self.send_error_response(400, "Content-Type must be multipart/form-data")
                    return
                
                # Extract boundary
                boundary = content_type.split("boundary=")[1].encode()
                content_length = int(self.headers.get('Content-Length', 0))
                body = self.rfile.read(content_length)
                
                # Split body by boundary to find the file and parameters
                parts = body.split(b"--" + boundary)
                file_data = None
                filename = "document.pdf"
                project_name = ""
                doc_type = "presales"
                
                for part in parts:
                    if b"name=\"project\"" in part:
                        part_lines = part.split(b"\r\n")
                        for line in part_lines:
                            line_stripped = line.strip()
                            if line_stripped and not line_stripped.startswith(b"Content-"):
                                project_name = line_stripped.decode('utf-8')
                                break
                    elif b"name=\"doc_type\"" in part:
                        part_lines = part.split(b"\r\n")
                        for line in part_lines:
                            line_stripped = line.strip()
                            if line_stripped and not line_stripped.startswith(b"Content-"):
                                doc_type = line_stripped.decode('utf-8')
                                break
                    elif b"filename=" in part:
                        header_end = part.find(b"\r\n\r\n")
                        fn_match = re.search(b'filename="([^"]+)"', part)
                        if fn_match:
                            filename = fn_match.group(1).decode('utf-8')
                        if header_end != -1:
                            file_data = part[header_end+4:-2]
                
                if not file_data:
                    self.send_error_response(400, "No file uploaded")
                    return
                
                temp_path = "temp_scope" + os.path.splitext(filename)[1].lower()
                with open(temp_path, "wb") as f:
                    f.write(file_data)
                
                scope_fields = {}
                components = []
                if doc_type == "agreed":
                    scope_fields = extract_agreed_document_details(temp_path, filename)
                else:
                    scope_fields, components = extract_scope_and_components(temp_path, filename)
                
                scope_fields["Project"] = project_name
                
                if os.path.exists(temp_path):
                    os.remove(temp_path)
                
                self.send_json_response(200, {
                    "status": "success",
                    "scope": scope_fields,
                    "components": components,
                    "doc_type": doc_type
                })
            except Exception as e:
                import traceback
                traceback.print_exc()
                self.send_error_response(500, f"Error parsing document: {str(e)}")
            return

        elif path == "/api/plan/upload":
            try:
                content_type = self.headers.get('Content-Type', '')
                if not content_type.startswith('multipart/form-data'):
                    self.send_error_response(400, "Content-Type must be multipart/form-data")
                    return
                
                boundary = content_type.split("boundary=")[1].encode()
                content_length = int(self.headers.get('Content-Length', 0))
                body = self.rfile.read(content_length)
                
                parts = body.split(b"--" + boundary)
                file_data = None
                for part in parts:
                    if b"filename=" in part:
                        header_end = part.find(b"\r\n\r\n")
                        if header_end != -1:
                            file_data = part[header_end+4:-2]
                            break
                
                if not file_data:
                    self.send_error_response(400, "No file uploaded")
                    return
                
                temp_pdf_path = "temp_plan.pdf"
                with open(temp_pdf_path, "wb") as f:
                    f.write(file_data)
                
                plan_data = parse_pdf_plan_with_hierarchy(temp_pdf_path)
                
                with open("project_plan.json", "w", encoding="utf-8") as f:
                    json.dump(plan_data, f, indent=2)
                
                if os.path.exists(temp_pdf_path):
                    os.remove(temp_pdf_path)
                
                self.send_json_response(200, {
                    "status": "success",
                    "message": f"Successfully parsed {len(plan_data['tasks'])} tasks.",
                    "projectName": plan_data["projectName"],
                    "tasks": plan_data["tasks"]
                })
            except Exception as e:
                import traceback
                traceback.print_exc()
                self.send_error_response(500, f"Error parsing PDF: {str(e)}")
            return

        elif path == "/api/plan/save":
            try:
                content_length = int(self.headers['Content-Length'])
                post_data = self.rfile.read(content_length)
                plan_data = json.loads(post_data.decode('utf-8'))
                
                with open("project_plan.json", "w", encoding="utf-8") as f:
                    json.dump(plan_data, f, indent=2)
                
                self.send_json_response(200, {
                    "status": "success",
                    "message": "Project plan updates saved permanently on server."
                })
            except Exception as e:
                import traceback
                traceback.print_exc()
                self.send_error_response(500, f"Error saving project plan: {str(e)}")
            return

        elif path == "/api/plan/delete":
            try:
                plan_path = "project_plan.json"
                if os.path.exists(plan_path):
                    os.remove(plan_path)
                self.send_json_response(200, {
                    "status": "success",
                    "message": "Project plan deleted successfully."
                })
            except Exception as e:
                import traceback
                traceback.print_exc()
                self.send_error_response(500, f"Error deleting project plan: {str(e)}")
            return

        elif path == "/api/approve":
            try:
                content_length = int(self.headers['Content-Length'])
                post_data = self.rfile.read(content_length)
                payload = json.loads(post_data.decode('utf-8'))

                project_name = payload.get("project")
                approved_by = payload.get("approvedBy")

                if not project_name:
                    self.send_error_response(400, "Missing project name.")
                    return

                # 1. Update database
                supabase.table("projects").update({"approved_by": approved_by}).eq("project", project_name).execute()

                # 2. Update Excel local cache
                if os.path.exists(EXCEL_PATH):
                    with EXCEL_LOCK:
                        excel_file = pd.ExcelFile(EXCEL_PATH)
                        sheets = {s: pd.read_excel(excel_file, sheet_name=s) for s in excel_file.sheet_names}
                        df_projects = sheets["Projects"]
                        if "Approved By" not in df_projects.columns:
                            df_projects["Approved By"] = None
                        df_projects["Approved By"] = df_projects["Approved By"].astype(object)

                        idx_list = df_projects.index[df_projects["Project"] == project_name].tolist()
                        if idx_list:
                            df_projects.at[idx_list[0], "Approved By"] = approved_by

                        with pd.ExcelWriter(EXCEL_PATH, engine="openpyxl") as writer:
                            for sheet_name, df_sheet in sheets.items():
                                df_sheet.to_excel(writer, sheet_name=sheet_name, index=False)

                self.send_json_response(200, {
                    "status": "success",
                    "message": f"Project approval status updated for {project_name}."
                })
            except Exception as e:
                import traceback
                traceback.print_exc()
                self.send_error_response(500, f"Error processing approval: {str(e)}")
            return

        elif path == "/api/approve/bulk":
            try:
                content_length = int(self.headers['Content-Length'])
                post_data = self.rfile.read(content_length)
                payload = json.loads(post_data.decode('utf-8'))

                approvals = payload.get("approvals")

                if approvals is None or not isinstance(approvals, dict):
                    self.send_error_response(400, "Missing approvals mapping.")
                    return

                # 1. Update database bulk
                for proj_name, approved_by in approvals.items():
                    supabase.table("projects").update({"approved_by": approved_by}).eq("project", proj_name).execute()

                # 2. Update Excel cache
                if os.path.exists(EXCEL_PATH):
                    with EXCEL_LOCK:
                        excel_file = pd.ExcelFile(EXCEL_PATH)
                        sheets = {s: pd.read_excel(excel_file, sheet_name=s) for s in excel_file.sheet_names}
                        df_projects = sheets["Projects"]
                        if "Approved By" not in df_projects.columns:
                            df_projects["Approved By"] = None
                        df_projects["Approved By"] = df_projects["Approved By"].astype(object)

                        for proj_name, approved_by in approvals.items():
                            idx_list = df_projects.index[df_projects["Project"] == proj_name].tolist()
                            if idx_list:
                                df_projects.at[idx_list[0], "Approved By"] = approved_by

                        with pd.ExcelWriter(EXCEL_PATH, engine="openpyxl") as writer:
                            for sheet_name, df_sheet in sheets.items():
                                df_sheet.to_excel(writer, sheet_name=sheet_name, index=False)

                self.send_json_response(200, {
                    "status": "success",
                    "message": "Bulk approval status updated."
                })
            except Exception as e:
                import traceback
                traceback.print_exc()
                self.send_error_response(500, f"Error in bulk approval: {str(e)}")
            return

        elif path == "/api/approve/reset":
            try:
                # 1. Reset database
                supabase.table("projects").update({"approved_by": None}).execute()
                supabase.table("approvals").delete().neq("approved_by", "dummy_unlikely_match_val").execute()

                # 2. Reset local Excel
                if os.path.exists(EXCEL_PATH):
                    with EXCEL_LOCK:
                        excel_file = pd.ExcelFile(EXCEL_PATH)
                        sheets = {s: pd.read_excel(excel_file, sheet_name=s) for s in excel_file.sheet_names}
                        df_projects = sheets["Projects"]
                        df_projects["Approved By"] = None
                        
                        df_approvals = sheets["Approvals"]
                        df_approvals = df_approvals.iloc[0:0] # clear log
                        sheets["Approvals"] = df_approvals

                        with pd.ExcelWriter(EXCEL_PATH, engine="openpyxl") as writer:
                            for sheet_name, df_sheet in sheets.items():
                                df_sheet.to_excel(writer, sheet_name=sheet_name, index=False)

                self.send_json_response(200, {
                    "status": "success",
                    "message": "All project approvals reset successfully."
                })
            except Exception as e:
                import traceback
                traceback.print_exc()
                self.send_error_response(500, f"Error resetting approvals: {str(e)}")
            return

        elif path == "/api/ppt/generate":
            try:
                content_length = int(self.headers['Content-Length'])
                post_data = self.rfile.read(content_length)
                payload = json.loads(post_data.decode('utf-8'))
                
                template_file = "Liv_ing_Template.pptx"
                output_file = "Liv_ing_Project_Weekly_Review.pptx"
                
                if not os.path.exists(template_file):
                    self.send_error_response(404, f"PowerPoint template {template_file} not found.")
                    return
                
                # Import python-pptx and fill template
                from pptx import Presentation
                prs = Presentation(template_file)
                
                # Slide 1: Review Date
                slide1 = prs.slides[0]
                for shape in slide1.shapes:
                    if shape.name == "TextBox 14":
                        shape.text_frame.text = payload.get("reviewDate", "")
                
                # Slide 2: Weekly Review Data
                slide2 = prs.slides[1]
                for shape in slide2.shapes:
                    if shape.name == "Rectangle: Rounded Corners 3":
                        shape.text_frame.text = payload.get("projectManager", "")
                    elif shape.name == "Rectangle: Rounded Corners 21":
                        shape.text_frame.text = payload.get("latestUpdates", "")
                    elif shape.name == "Rectangle: Rounded Corners 12":
                        shape.text_frame.text = payload.get("risks", "")
                    elif shape.name == "Table 23":
                        table = shape.table
                        milestones = payload.get("milestones", [])
                        for idx, row_data in enumerate(milestones):
                            r_idx = idx + 2
                            if r_idx < len(table.rows):
                                table.cell(r_idx, 0).text = row_data.get("Activity", "")
                                table.cell(r_idx, 1).text = row_data.get("ETA", "")
                                table.cell(r_idx, 2).text = row_data.get("Rev", "")
                                table.cell(r_idx, 3).text = row_data.get("UpdatedPlan", "")
                                table.cell(r_idx, 4).text = row_data.get("Status", "")
                                table.cell(r_idx, 5).text = row_data.get("Comments", "")
                
                # Slide 3: Support Required
                slide3 = prs.slides[2]
                for shape in slide3.shapes:
                    if shape.name == "Content Placeholder 2":
                        shape.text_frame.text = payload.get("supportRequired", "")
                
                prs.save(output_file)
                
                # Read file bytes
                with open(output_file, "rb") as f:
                    file_bytes = f.read()
                
                self.send_response(200)
                self.send_header("Content-Type", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
                self.send_header("Content-Disposition", "attachment; filename=Liv_ing_Project_Weekly_Review.pptx")
                self.send_header("Content-Length", str(len(file_bytes)))
                self.end_headers()
                self.wfile.write(file_bytes)
                
                # Clean up local output file
                if os.path.exists(output_file):
                    os.remove(output_file)
                
            except Exception as e:
                import traceback
                traceback.print_exc()
                self.send_error_response(500, f"Error generating PowerPoint: {str(e)}")
            return

        self.send_error_response(404, "Endpoint not found.")

    def send_json_response(self, status_code, data):
        try:
            sanitized = sanitize_nan(data)
            response_body = json.dumps(sanitized, allow_nan=False).encode('utf-8')
        except ValueError as ve:
            print(f"JSON serialization error (NaN/Inf detected): {ve}")
            self.send_error_response(500, "Internal server error: Invalid float values (NaN/Infinity) encountered.")
            return

        self.send_response(status_code)
        self.send_header("Content-Type", "application/json")
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Content-Length", str(len(response_body)))
        self.end_headers()
        self.wfile.write(response_body)

    def send_error_response(self, status_code, message):
        try:
            sanitized = sanitize_nan({"status": "error", "message": message})
            response_body = json.dumps(sanitized, allow_nan=False).encode('utf-8')
        except ValueError as ve:
            print(f"JSON serialization error in error response: {ve}")
            response_body = b'{"status": "error", "message": "Internal server error"}'
        self.send_response(status_code)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(response_body)))
        self.end_headers()
        self.wfile.write(response_body)

def run_server():
    server_address = ('', PORT)
    httpd = HTTPServer(server_address, PMORequestHandler)
    print(f"PMO Web Portal local server running at: http://localhost:{PORT}")
    try:
        httpd.serve_forever()
    except KeyboardInterrupt:
        print("\nShutting down server.")
        httpd.server_close()

if __name__ == "__main__":
    run_server()
